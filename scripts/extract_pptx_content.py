"""Extract source PPTX text to UTF-8 file."""
import json
from pathlib import Path
from pptx import Presentation

downloads = Path.home() / "Downloads"
source = None
for p in downloads.glob("*.pptx"):
    if "AJITEU" in p.name.upper() or "아지트" in p.name or "포트폴리오" in p.name:
        source = p
        break

if not source:
    source = Path(r"c:\Users\PC\Downloads\팀 프로젝트 '아지트(AJITEU)' 1포트폴리오.pptx")

prs = Presentation(str(source))
slides = []
for i, slide in enumerate(prs.slides, 1):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    slides.append({"index": i, "texts": texts})

out = Path(__file__).resolve().parents[1] / "scripts" / "source_ppt_content.json"
out.write_text(json.dumps({"source": str(source), "slides": slides}, ensure_ascii=False, indent=2), encoding="utf-8")
print(f"Saved {len(slides)} slides to {out}")
