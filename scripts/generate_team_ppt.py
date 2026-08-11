"""Ajiteu 팀 발표용 PPT 생성 스크립트."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = Path(__file__).resolve().parent / "ppt_assets"
OUTPUT_KO = ROOT / "Ajiteu_팀발표.pptx"
OUTPUT_EN = ROOT / "Ajiteu_Team_Presentation.pptx"

PRIMARY = RGBColor(0x93, 0x81, 0xD6)
DARK = RGBColor(0x26, 0x26, 0x26)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF3, 0xFF)

PROJECT_PITCH = (
    "자신의 일상을 사진·글로 자유롭게 기록하고 공유하며, "
    "좋아요·댓글을 통해 사람들과 소통하고 "
    "새로운 관심사와 인연을 발견할 수 있는 라이프스타일 커뮤니티입니다."
)


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, title: str, subtitle: str = "") -> None:
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.5), Inches(0.72), Inches(9), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = RGBColor(0xEE, 0xEA, 0xFF)


def add_bullets(slide, items: list[str], left=0.7, top=1.5, width=8.8, height=5.5, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)


def add_two_columns(slide, left_title, left_items, right_title, right_items, size=15):
    lt = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(4.3), Inches(0.45))
    lp = lt.text_frame.paragraphs[0]
    lp.text = left_title
    lp.font.bold = True
    lp.font.size = Pt(17)
    lp.font.color.rgb = PRIMARY

    lb = slide.shapes.add_textbox(Inches(0.55), Inches(1.85), Inches(4.3), Inches(5.0))
    ltf = lb.text_frame
    ltf.word_wrap = True
    for i, item in enumerate(left_items):
        p = ltf.paragraphs[0] if i == 0 else ltf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(5)

    rt = slide.shapes.add_textbox(Inches(5.15), Inches(1.35), Inches(4.3), Inches(0.45))
    rp = rt.text_frame.paragraphs[0]
    rp.text = right_title
    rp.font.bold = True
    rp.font.size = Pt(17)
    rp.font.color.rgb = PRIMARY

    rb = slide.shapes.add_textbox(Inches(5.15), Inches(1.85), Inches(4.3), Inches(5.0))
    rtf = rb.text_frame
    rtf.word_wrap = True
    for i, item in enumerate(right_items):
        p = rtf.paragraphs[0] if i == 0 else rtf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.space_after = Pt(5)


def add_screenshot_slide(
    prs: Presentation,
    title: str,
    subtitle: str,
    image_path: Path,
    caption: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, title, subtitle)

    if image_path.exists():
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.55),
            Inches(1.25),
            width=Inches(8.9),
        )

    if caption:
        cap = slide.shapes.add_textbox(Inches(0.55), Inches(6.85), Inches(8.9), Inches(0.5))
        cp = cap.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(12)
        cp.font.color.rgb = GRAY
        cp.alignment = PP_ALIGN.CENTER


def slide_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, LIGHT_BG)

    accent = slide.shapes.add_shape(1, Inches(0), Inches(2.2), Inches(10), Inches(2.6))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.45), Inches(8.5), Inches(1.0))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Ajiteu (아지트)"
    tp.font.size = Pt(44)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.45), Inches(8.5), Inches(0.6))
    sp = sub.text_frame.paragraphs[0]
    sp.text = "함께 나누는 일상, 더 가까워지는 우리"
    sp.font.size = Pt(20)
    sp.font.color.rgb = RGBColor(0xEE, 0xEA, 0xFF)

    pitch = slide.shapes.add_textbox(Inches(0.8), Inches(4.15), Inches(8.5), Inches(1.2))
    pp = pitch.text_frame
    pp.word_wrap = True
    pp.paragraphs[0].text = PROJECT_PITCH
    pp.paragraphs[0].font.size = Pt(14)
    pp.paragraphs[0].font.color.rgb = RGBColor(0xEE, 0xEA, 0xFF)

    info = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(8.5), Inches(1.2))
    tf = info.text_frame
    for i, line in enumerate([
        "팀 프로젝트 발표  |  Flask · Python · SQLite · Bootstrap 5",
        "GitHub: github.com/t01040588614-source/Ajiteu",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = GRAY


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)

    # 목차
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "목차", "Agenda")
    add_bullets(slide, [
        "1. 프로젝트 소개",
        "2. 팀 구성 및 역할",
        "3. 화면 데모 (로그인 · 메인 · 모달)",
        "4. 기술 스택 & 주요 기능",
        "5. 시스템 아키텍처 & DB",
        "6. 통합 과정 & 트러블슈팅",
        "7. 잘한 점 & 아쉬운 점",
        "8. 실행 방법 & Q&A",
    ], size=19)

    # 프로젝트 소개
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "프로젝트 소개", "Lifestyle Community SNS")
    add_bullets(slide, [
        "일상을 사진·글로 자유롭게 기록하고 공유하는 라이프스타일 커뮤니티",
        "여행, 음식, 운동, 공부, 반려동물, 취미 등 다양한 주제 공유",
        "좋아요·댓글·답글로 소통, 카테고리·트렌드로 관심사 발견",
        "페이지 이동을 줄인 overlay 모달 UX로 가볍게 탐색",
        "",
        "▶ 한 줄 소개",
        PROJECT_PITCH,
    ], size=16)

    # 팀 구성 A/B
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "팀 구성 및 역할 (1/2)", "Team Roles")
    add_two_columns(
        slide,
        "👤 A — 성경수 | 인증·프로필",
        [
            "로그인 / 회원가입 / 로그아웃",
            "세션 기반 인증 (Flask g.user)",
            "프로필 조회·수정 (닉네임, 소개, 사진)",
            "auth_views.py, profile.py",
            "초기: 개별 DB → 통합 시 스키마 맞춤",
        ],
        "👤 B — 김국진 | 메인 UI",
        [
            "메인 화면 3단 레이아웃 (좌·중·우)",
            "게시글 목록 2열 카드 그리드",
            "반응형 CSS, 검색·카테고리 메뉴",
            "프로필 사이드바, 게시글 카드 UI",
            "main.css, components.css, pagination.css",
        ],
        size=14,
    )

    # 팀 구성 C/D
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "팀 구성 및 역할 (2/2)", "Team Roles")
    add_two_columns(
        slide,
        "👤 C — 김시내 | 게시글·댓글",
        [
            "글쓰기 / 게시글 수정·삭제 / 상세보기",
            "다중 이미지 업로드·갤러리(Carousel)",
            "댓글·답글 작성 / 수정 / 삭제",
            "post_api, comment_api, reply_api",
            "overlay 모달 내 댓글 UX 구현",
        ],
        "👤 D — 김환석 | DB·통합·QA",
        [
            "SQLite + SQLAlchemy ORM 연동",
            "좋아요·좋아요 수·조회수 API",
            "Blueprint 통합, 마이그레이션",
            "최종 통합·버그 수정·README·PPT",
            "fetch overlay JS, event delegation",
        ],
        size=14,
    )

    # 화면 데모 — 스크린샷
    add_screenshot_slide(
        prs,
        "화면 데모 ① 로그인",
        "담당: 성경수",
        ASSETS / "01_login.png",
        "그라데이션 브랜드 영역 + 로그인·회원가입 폼 — 서비스 첫인상",
    )
    add_screenshot_slide(
        prs,
        "화면 데모 ② 메인 피드",
        "담당: 김국진",
        ASSETS / "02_main.png",
        "3단 레이아웃 · 2열 카드 · 검색 · 카테고리 · 이번주의 발견",
    )
    add_screenshot_slide(
        prs,
        "화면 데모 ③ 상세 모달",
        "담당: 김시내 · 김환석",
        ASSETS / "03_modal.png",
        "overlay 상세 · 좋아요/댓글/조회수 · ⋯ 액션 메뉴 · 댓글 입력",
    )

    # 기술 스택
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "기술 스택", "Tech Stack")
    add_two_columns(
        slide,
        "Backend",
        [
            "Python 3.x / Flask 3.1 (Blueprint)",
            "Flask-SQLAlchemy · Flask-Migrate",
            "Flask-WTF / WTForms · Werkzeug",
            "SQLite (로컬) · Jinja2",
        ],
        "Frontend & 협업",
        [
            "Bootstrap 5.3 · Vanilla JS (fetch)",
            "CSS 모듈화 (main/detail/components)",
            "Git / GitHub · README 문서화",
            "python-pptx (발표 자료 생성)",
        ],
        size=16,
    )

    # 주요 기능
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "주요 기능 & UI", "Key Features")
    add_two_columns(
        slide,
        "사용자·게시물",
        [
            "회원가입 / 로그인 / 프로필 설정",
            "게시글 CRUD + 다중 이미지",
            "반투명 overlay 상세 모달",
            "⋯ 액션 메뉴 (수정/삭제/이동)",
            "이번주의 발견 · 카테고리 필터",
        ],
        "소셜·인터랙션",
        [
            "댓글 / 답글 CRUD",
            "게시물 좋아요 (AJAX JSON)",
            "조회수 (상세 열람 시에만 +1)",
            "키워드 검색 · 페이지네이션",
            "event delegation 클릭 처리",
        ],
        size=14,
    )

    # 아키텍처
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "시스템 아키텍처", "Architecture")
    add_bullets(slide, [
        "[Browser] HTML/CSS/JS  ←→  fetch · form submit",
        "[Flask] app.py → create_app() → Blueprint",
        "  auth_views  post_api  comment_api  reply_api  profile",
        "[DB] SQLite ←→ User · Post · Comment · Reply · post_liker",
        "[패턴] 클릭 → API → 성공 → 재조회 → render (단일 데이터 흐름)",
    ], size=17)

    # DB
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "데이터베이스 설계", "Database")
    add_two_columns(
        slide,
        "테이블",
        [
            "user — username, nickname, intro, image",
            "post — content, image_path, view_count",
            "comment / reply — content, user_id",
            "post_liker (M:N) — 좋아요 관계",
        ],
        "통합 과정",
        [
            "초기: 팀원별 개별 DB 개발",
            "ORM 모델·관계 통일 (1:N, M:N)",
            "Flask-Migrate로 스키마 버전 관리",
            "단일 SQLite로 최종 통합",
        ],
        size=15,
    )

    # 트러블슈팅
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "통합 & 트러블슈팅", "Issues & Solutions")
    add_two_columns(
        slide,
        "발생한 문제",
        [
            "팀원별 DB·API 불일치 → 통합 난이도",
            "댓글 DOM 직접 수정 → 화면·DB 불일치",
            "render마다 이벤트 리스너 중복",
            "조회수: 댓글 갱신·수정 시에도 +1",
            "overlay 잔류, 1번째 댓글 수정 무반응",
            "이번주의 발견 클릭 미동작",
        ],
        "해결 방법",
        [
            "단일 SQLite + Blueprint 통합",
            "서버 재조회 → comments render 패턴",
            "modal 1개 click/submit delegation",
            "track_view=0 으로 조회수 분리",
            "data-action + closest() 버튼 통일",
            "weeklyDiscoveryList delegation",
        ],
        size=13,
    )

    # 잘한 점 / 아쉬운 점
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "잘한 점 & 아쉬운 점", "Retrospective")
    add_two_columns(
        slide,
        "✅ 잘 표현·개선한 부분",
        [
            "Blueprint로 기능별 모듈 분리",
            "overlay 모달로 SPA-like UX 구현",
            "fetch + partial render로 모달 유지",
            "event delegation으로 클릭 안정화",
            "조회수·댓글 상태 동기화 개선",
            "팀 GitHub 통합·README 문서화",
        ],
        "⚠️ 실수·아쉬운 부분",
        [
            "초기 각자 DB → 통합 비용 발생",
            "DOM 직접 조작으로 버그 누적",
            "inline onclick 혼용 → 이벤트 충돌",
            "조회수 로직을 detail API에만 의존",
            "통합 테스트·자동화 테스트 부족",
            "팔로우·알림 등 확장 기능 미구현",
        ],
        size=13,
    )

    # 실행 방법
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_title_bar(slide, "실행 방법 & 데모", "How to Run")
    add_bullets(slide, [
        "git clone https://github.com/t01040588614-source/Ajiteu.git",
        "python -m venv venv  →  pip install -r requirements.txt",
        "config.py (SECRET_KEY, SQLALCHEMY_DATABASE_URI)",
        "flask db upgrade  →  python app.py",
        "http://127.0.0.1:5000",
        "",
        "데모: 로그인 → 글쓰기 → 상세 모달 → 댓글·좋아요 → 프로필",
    ], size=17)

    # Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)
    qa = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    qp = qa.text_frame.paragraphs[0]
    qp.text = "Q & A"
    qp.font.size = Pt(52)
    qp.font.bold = True
    qp.font.color.rgb = WHITE
    qp.alignment = PP_ALIGN.CENTER

    thanks = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(8.4), Inches(1.5))
    tf = thanks.text_frame
    tf.word_wrap = True
    for i, line in enumerate([
        "경청해 주셔서 감사합니다",
        PROJECT_PITCH,
        "GitHub: github.com/t01040588614-source/Ajiteu",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16 if i else 22)
        p.font.color.rgb = RGBColor(0xEE, 0xEA, 0xFF)
        p.alignment = PP_ALIGN.CENTER

    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT_KO)
    prs.save(OUTPUT_EN)
    print(f"Created: {OUTPUT_KO}")
    print(f"Created: {OUTPUT_EN}")


if __name__ == "__main__":
    main()
