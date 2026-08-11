"""포트폴리오 PPTX를 바탕으로 팀 발표용 PPT 생성."""

import json
import shutil
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
META = json.loads((Path(__file__).parent / "source_ppt_content.json").read_text(encoding="utf-8"))
SOURCE = Path(META["source"])
OUTPUT_KO = ROOT / "Ajiteu_팀발표.pptx"
OUTPUT_EN = ROOT / "Ajiteu_Team_Presentation.pptx"

# 원본 텍스트 → 최신 프로젝트 반영 텍스트
TEXT_REPLACEMENTS = {
    "Flask & Python 기반의 사용자 중심 커뮤니티 및 SNS 웹 애플리케이션 개발 프로젝트":
        "Flask & Python 기반 커뮤니티 SNS | GitHub: github.com/t01040588614-source/Ajiteu",
    "중앙 메인 피드: 글 작성/수정, 이미지 피드, 좋아요/댓글/조회수 실시간 카드 표시":
        "중앙 메인 피드: 2열 카드 그리드, 반투명 overlay 모달 상세보기, 좋아요/댓글/조회수",
    "게시글 & 미디어 피드":
        "게시글 & Overlay 모달 UX",
    "이미지 첨부 피드 작성, 상세 보기, 좋아요 기능, 비동기 댓글/대댓글 작성 및 실시간 조회수 카운팅을 지원합니다.":
        "fetch 기반 overlay 모달로 게시물·프로필·글쓰기 팝업, 다중 이미지 갤러리(Carousel), "
        "댓글/답글 작성·수정·삭제, ⋯ 액션 메뉴, 조회수 카운팅을 지원합니다.",
    "comment_api.py / reply_api.py : 댓글 및 대댓글 AJAX API":
        "comment_api.py / reply_api.py : 댓글·답글 CRUD API",
    "post_api.py : 게시글 생성, 수정, 삭제 비동기 처리":
        "post_api.py : 게시글 CRUD, 좋아요(JSON), 검색·카테고리 필터",
    "컴포넌트화된 CSS와 모듈화된 JS를 적용하여 메인 피드와 프로필 관리의 UI/UX 반응성을 대폭 향상시켰습니다.":
        "openPostDetailModal(), openProfileModal(), openWriteModal() 등 overlay JS와 "
        "detail.css/components.css로 모달·갤러리·액션시트 UI를 구현했습니다.",
    "비동기 API 처리: 댓글 작성 및 좋아요 클릭 시 불필요한 전체 페이지 새로고침 방지":
        "Overlay fetch 처리: 댓글 등록·삭제 시 모달 유지 + 영역만 갱신, 좋아요 AJAX",
    "기능 단위 Blueprint 분리와 RESTful API 도입을 통해 코드의 재사용성을 높이고 유지보수 비용을 크게 절감했습니다.":
        "Blueprint 모듈 분리, overlay 모달 UX, 댓글 삭제/수정 버그 해결로 "
        "사용자 경험과 코드 유지보수성을 동시에 개선했습니다.",
    "3. 프론트 연동 & UI":
        "3. 프론트 연동 & Overlay UI",
    "Jinja2 템플릿 완성, 모듈별 CSS 스타일링 및 JS 비동기 통신 구현":
        "Jinja2 템플릿, CSS 모듈화, fetch overlay JS, 댓글/모달 이벤트 위임 처리",
    "4. 테스트 & 디버깅":
        "4. 테스트 & 모달 디버깅",
    "SQLite DB 통합 테스트, 예외 처리 강화 및 최종 배포 환경 설정":
        "모달 닫기·overlay·댓글 수정/삭제 QA, fetch 실패 예외 처리, README 문서화",
    "프로젝트: 아지트(AJITEU)  Tech: Flask, Python, SQLite":
        "프로젝트: 아지트(AJITEU)  |  Tech: Flask, Python, SQLite, Bootstrap 5",
}

# 슬라이드별 shape index → 새 텍스트 (전체 교체)
SLIDE_TEXT_OVERRIDES = {
    1: {
        2: "팀 프로젝트 '아지트(AJITEU)'\n 팀 발표",
    },
    12: {
        2: "경청해 주셔서 감사합니다.\nGitHub: github.com/t01040588614-source/Ajiteu",
    },
}


def set_shape_text(shape, new_text: str) -> None:
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = new_text


def apply_replacements(prs: Presentation) -> None:
    for slide_idx, slide in enumerate(prs.slides, 1):
        overrides = SLIDE_TEXT_OVERRIDES.get(slide_idx, {})
        for shape_idx, shape in enumerate(slide.shapes):
            if not hasattr(shape, "text") or not shape.text.strip():
                continue

            if shape_idx in overrides:
                set_shape_text(shape, overrides[shape_idx])
                continue

            original = shape.text
            updated = original
            for old, new in TEXT_REPLACEMENTS.items():
                if old in updated:
                    updated = updated.replace(old, new)
            if updated != original:
                set_shape_text(shape, updated)


def main() -> None:
    if not SOURCE.exists():
        raise FileNotFoundError(f"원본 PPT를 찾을 수 없습니다: {SOURCE}")

    for output in (OUTPUT_KO, OUTPUT_EN):
        shutil.copy2(SOURCE, output)
        prs = Presentation(str(output))
        apply_replacements(prs)
        prs.save(str(output))
        print(f"Created: {output}")


if __name__ == "__main__":
    main()
