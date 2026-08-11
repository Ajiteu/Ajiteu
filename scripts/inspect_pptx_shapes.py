"""Inspect shape details in source PPTX."""
import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

source = json.loads((Path(__file__).parent / "source_ppt_content.json").read_text(encoding="utf-8"))["source"]
prs = Presentation(source)

report = []
for si, slide in enumerate(prs.slides, 1):
    slide_info = {"slide": si, "shapes": []}
    for idx, shape in enumerate(slide.shapes):
        info = {
            "idx": idx,
            "name": shape.name,
            "type": str(shape.shape_type),
            "has_text": hasattr(shape, "text"),
        }
        if hasattr(shape, "text"):
            info["text"] = shape.text
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            info["group_shapes"] = len(shape.shapes)
        slide_info["shapes"].append(info)
    report.append(slide_info)

out = Path(__file__).parent / "source_ppt_shapes.json"
out.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
print("saved", out)
