"""Inspect source PPTX file."""
from pathlib import Path
from pptx import Presentation

candidates = [
    Path(r"c:\Users\PC\Downloads\팀 프로젝트 '아지트(AJITEU)' 1포트폴리오.pptx"),
    Path.home() / "Downloads",
]

downloads = Path.home() / "Downloads"
if downloads.exists():
    for p in downloads.glob("*.pptx"):
        print("FOUND:", p)
        candidates.insert(0, p)

for path in candidates:
    if isinstance(path, Path) and path.is_file() and path.suffix == ".pptx":
        print("\n=== Analyzing:", path, "===")
        prs = Presentation(str(path))
        print("slides:", len(prs.slides))
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            print(f"\n--- Slide {i} ---")
            for t in texts:
                print(t[:200])
